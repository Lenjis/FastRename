from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

# 创建PPT
presentation = Presentation()

# 配色
xmu_blue = RGBColor(0, 56, 168)
xmu_gray = RGBColor(220, 220, 220)
white = RGBColor(255, 255, 255)

# 字体大小
title_font_size = Pt(40)
subtitle_font_size = Pt(24)
content_font_size = Pt(18)

# 添加首页
slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # 空白布局
# 添加标题框
title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "马克思主义基本原理\n课程汇报"
title_frame.paragraphs[0].font.size = title_font_size
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = xmu_blue
title_frame.paragraphs[0].alignment = 1  # 居中对齐
# 添加副标题框
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "厦门大学 - 汇报人：XXX - 日期：XXXX年XX月"
subtitle_frame.paragraphs[0].font.size = subtitle_font_size
subtitle_frame.paragraphs[0].font.color.rgb = xmu_blue
subtitle_frame.paragraphs[0].alignment = 1

# 添加问题引入页
slide = presentation.slides.add_slide(presentation.slide_layouts[1])  # 带标题和内容布局
title = slide.shapes.title
title.text = "问题引入"
title.text_frame.paragraphs[0].font.color.rgb = xmu_blue
title.text_frame.paragraphs[0].font.size = title_font_size
content = slide.placeholders[1]
content.text = "1. 背景说明\n2. 问题描述"
content.text_frame.paragraphs[0].font.size = content_font_size

# 添加机制分析页
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "机制分析"
title.text_frame.paragraphs[0].font.color.rgb = xmu_blue
title.text_frame.paragraphs[0].font.size = title_font_size
content = slide.placeholders[1]
content.text = "1. 理论框架\n2. 模型图示"
content.text_frame.paragraphs[0].font.size = content_font_size

# 添加问卷设计页
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "问卷设计"
title.text_frame.paragraphs[0].font.color.rgb = xmu_blue
title.text_frame.paragraphs[0].font.size = title_font_size
content = slide.placeholders[1]
content.text = "1. 设计流程\n2. 问卷示例"
content.text_frame.paragraphs[0].font.size = content_font_size

# 添加结果分析页
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "结果分析"
title.text_frame.paragraphs[0].font.color.rgb = xmu_blue
title.text_frame.paragraphs[0].font.size = title_font_size
content = slide.placeholders[1]
content.text = "1. 数据图表\n2. 关键结论"
content.text_frame.paragraphs[0].font.size = content_font_size

# 添加解决方案页
slide = presentation.slides.add_slide(presentation.slide_layouts[1])
title = slide.shapes.title
title.text = "解决方案"
title.text_frame.paragraphs[0].font.color.rgb = xmu_blue
title.text_frame.paragraphs[0].font.size = title_font_size
content = slide.placeholders[1]
content.text = "1. 对策与展望"
content.text_frame.paragraphs[0].font.size = content_font_size

# 添加结束页
slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # 空白布局
# 添加标题框
title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "感谢聆听"
title_frame.paragraphs[0].font.size = title_font_size
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = xmu_blue
title_frame.paragraphs[0].alignment = 1  # 居中对齐
# 添加副标题框
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "联系方式：xxxx@example.com"
subtitle_frame.paragraphs[0].font.size = subtitle_font_size
subtitle_frame.paragraphs[0].font.color.rgb = xmu_blue
subtitle_frame.paragraphs[0].alignment = 1

# 保存文件
file_path = "./马克思主义基本原理课程汇报模板.pptx"
presentation.save(file_path)

file_path
